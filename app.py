import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io, os, tempfile, re, time
from docx import Document
import google.generativeai as genai
from PIL import Image

# --- CONFIGURACIÓN DE PÁGINA ---
st.set_page_config(page_title="Generador Universal de Negocios", layout="wide")
st.title("🏆 Business Case Generator: Presentaciones de Inversión")
st.markdown("Analiza Lavaderos, Inmuebles, Apps o cualquier unidad de negocio con IA Multimodal.")

# --- BARRA LATERAL ---
with st.sidebar:
    st.header("⚙️ Configuración")
    api_key = st.text_input("Introduce tu Google API Key", type="password")
    moneda_local = st.text_input("Moneda Local (ej: UYU, ARS)", value="UYU")
    st.divider()
    st.info("KPIs incluidos: ROTE, ROI, Cost to Income, Punto de Equilibrio.")

# --- FUNCIONES TÉCNICAS BLINDADAS ---

def hex_to_rgb(hex_color):
    """Convierte HEX a RGB para PowerPoint"""
    try:
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    except:
        return RGBColor(212, 175, 55) # Dorado por defecto si falla

def leer_archivo_texto(file):
    """Evita el error de Zip diferenciando Word de TXT"""
    nombre = file.name.lower()
    try:
        if nombre.endswith('.docx'):
            doc = Document(file)
            return "\n".join([p.text for p in doc.paragraphs])
        else:
            return file.read().decode("utf-8", errors="ignore")
    except Exception as e:
        return f"No se pudo leer el archivo: {e}"

def disenar_slide(slide, titulo, contenido, accent_hex, imagen_pil=None):
    """Motor de Diseño: Fondo oscuro, títulos en color de acento y maquetación pro"""
    # Fondo Oscuro Premium
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(18, 18, 18)
    
    color_acento = hex_to_rgb(accent_hex)
    
    # Título
    title_shape = slide.shapes.title
    title_shape.text = titulo.upper()
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = color_acento
    title_para.alignment = PP_ALIGN.LEFT

    # Cuerpo de texto
    body_shape = slide.placeholders[1]
    # Si hay imagen, el texto ocupa la mitad
    body_shape.width = Inches(5.2) if imagen_pil else Inches(9)
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    for linea in contenido.split('\n'):
        clean_l = linea.strip().replace('*', '').replace('-', '').strip()
        if len(clean_l) > 3:
            p = tf.add_paragraph()
            p.text = "• " + clean_l
            p.font.size = Pt(17)
            p.font.color.rgb = RGBColor(230, 230, 230)
            p.space_after = Pt(10)

    # Insertar Imagen si existe (a la derecha)
    if imagen_pil:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            imagen_pil.save(tmp.name)
            slide.shapes.add_picture(tmp.name, Inches(5.6), Inches(1.5), height=Inches(4.5))

# --- INTERFAZ DE USUARIO (TODOS LOS BOTONES REQUERIDOS) ---
st.subheader("📥 Carga Multimodal de Proyecto")
c1, c2, c3 = st.columns(3)

with c1:
    f_doc = st.file_uploader("1. Conversación (Word o TXT)", type=["docx", "txt"])
    f_xlsx = st.file_uploader("2. Planilla Excel (Costos/Ventas)", type=["xlsx"])

with c2:
    f_media = st.file_uploader("3. Audio o Video de la charla", type=["mp3", "mp4", "wav", "m4a"])
    f_foto = st.file_uploader("4. Imagen (Fachada, Logo o Render)", type=["jpg", "png", "jpeg"])

with c3:
    f_notas = st.text_area("5. Notas adicionales y contexto:", height=160, placeholder="Explica el negocio aquí...")

# --- PROCESAMIENTO ---
if st.button("🚀 GENERAR PRESENTACIÓN COMERCIAL"):
    if not api_key:
        st.error("Introduce tu API Key en la barra lateral.")
    else:
        try:
            genai.configure(api_key=api_key)
            
            # --- AUTO-DETECCIÓN DE MODELO (EVITA ERROR 404) ---
            modelos_disponibles = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
            modelo_nombre = next((m for m in modelos_disponibles if '1.5-flash' in m), modelos_disponibles[0])
            model = genai.GenerativeModel(modelo_nombre)
            
            with st.spinner(f"Analizando negocio con {modelo_nombre}..."):
                # Unificar contexto
                contexto_final = f"MONEDA LOCAL: {moneda_local}\n"
                if f_notas: contexto_final += f"CONTEXTO: {f_notas}\n"
                if f_doc: contexto_final += f"DOCUMENTO: {leer_archivo_texto(f_doc)}\n"
                if f_xlsx:
                    df = pd.read_excel(f_xlsx)
                    contexto_final += f"DATOS EXCEL:\n{df.to_string()}\n"
                
                prompt_parts = [contexto_final]
                if f_foto: prompt_parts.append(Image.open(f_foto))
                
                # Manejo de Multimedia
                if f_media:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(f_media.name)[1]) as tmp:
                        tmp.write(f_media.read())
                        t_path = tmp.name
                    gf = genai.upload_file(path=t_path)
                    while gf.state.name == "PROCESSING":
                        time.sleep(2)
                        gf = genai.get_file(gf.name)
                    prompt_parts.append(gf)

                # PROMPT MAESTRO (UNIVERSAL)
                prompt_maestro = f"""
                Eres un Consultor Senior de Inversiones. Analiza el material adjunto.
                1. Detecta el sector de negocio.
                2. Elige un color HEX vibrante para el branding.
                3. Calcula KPIs: ROTE, ROI, Cost to Income, Punto de Equilibrio y otros datos importantes sugeridos.
                4. Explica flujos en {moneda_local} y USD.
                5. Menciona Beneficios Fiscales (Vivienda Promovida/COMAP si es Uruguay).
                
                RESPONDE CON ESTE FORMATO:
                COLOR: #HEX
                [SLIDE] Título | Punto 1 * Punto 2 * Punto 3
                
                Genera 10 diapositivas: Portada, Problema, Solución, Público, Ubicación/Zonas, Inversión, Operación, Beneficios Fiscales, KPIs Financieros, Cierre.
                """
                
                res = model.generate_content([prompt_maestro] + prompt_parts)
                texto_ia = res.text

                # --- CONSTRUCCIÓN DEL PPTX ---
                prs = Presentation()
                
                # Extraer color de acento
                match_color = re.search(r'COLOR:\s*(#[0-9A-Fa-f]{6})', texto_ia)
                accent_hex = match_color.group(1) if match_color else "#D4AF37"
                
                # Crear Slides
                bloques = texto_ia.split("[SLIDE]")
                for i, bloque in enumerate(bloques):
                    if "|" in bloque:
                        partes = bloque.split("|")
                        titulo_s = partes[0].replace(f"COLOR: {accent_hex}", "").strip()
                        contenido_s = partes[1].strip().replace("*", "\n")
                        
                        # Solo poner imagen en la portada o slides iniciales
                        img_para_slide = Image.open(f_foto) if (i <= 2 and f_foto) else None
                        disenar_slide(prs, titulo_s, contenido_s, accent_hex, img_para_slide)

                # Exportar para descarga
                buf = io.BytesIO()
                prs.save(buf)
                buf.seek(0)
                st.success("✅ ¡Presentación Comercial Generada!")
                st.download_button("📥 Descargar PowerPoint Profesional", buf, "Propuesta_Inversion.pptx")
                
                if f_media: os.remove(t_path)

        except Exception as e:
            st.error(f"Error técnico detectado: {e}")
